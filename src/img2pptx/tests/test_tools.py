"""Tests for img2pptx tools."""

import json
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from openai import AuthenticationError, OpenAIError, RateLimitError
from pptx import Presentation

from core import SHARED_WORKSPACE, get_workspace

from img2pptx.tools import (
    _create_slide,
    _extract_slide_content,
    _image_to_base64,
    _images_to_pptx_impl,
    _validate_image_path,
    _validate_output_path,
)


def test_validate_output_path_requires_pptx():
    """Output must have .pptx extension."""
    with pytest.raises(ValueError, match="must be .pptx"):
        _validate_output_path(Path("/tmp/test.pdf"))


def test_validate_output_path_forbids_system_dirs():
    """Cannot write to system directories."""
    with pytest.raises(ValueError, match="system directory"):
        _validate_output_path(Path("/etc/passwd.pptx"))


def test_validate_image_path_not_found():
    """Raises FileNotFoundError for missing images."""
    with pytest.raises(FileNotFoundError):
        _validate_image_path(Path("/nonexistent/image.png"))


def test_validate_image_path_bad_extension(tmp_path):
    """Rejects unsupported image formats."""
    bad_file = tmp_path / "test.txt"
    bad_file.write_text("not an image")
    with pytest.raises(ValueError, match="Unsupported format"):
        _validate_image_path(bad_file)


def test_get_workspace_path():
    """Workspace path should be valid directory."""
    path = get_workspace(SHARED_WORKSPACE)
    assert path.exists()
    assert path.is_dir()


def test_image_to_base64(tmp_path):
    """Should convert image to base64 data URL."""
    # Create a minimal PNG (1x1 transparent pixel)
    png_data = bytes(
        [
            0x89,
            0x50,
            0x4E,
            0x47,
            0x0D,
            0x0A,
            0x1A,
            0x0A,  # PNG signature
            0x00,
            0x00,
            0x00,
            0x0D,
            0x49,
            0x48,
            0x44,
            0x52,  # IHDR chunk
            0x00,
            0x00,
            0x00,
            0x01,
            0x00,
            0x00,
            0x00,
            0x01,
            0x08,
            0x06,
            0x00,
            0x00,
            0x00,
            0x1F,
            0x15,
            0xC4,
            0x89,
            0x00,
            0x00,
            0x00,
            0x0A,
            0x49,
            0x44,
            0x41,  # IDAT chunk
            0x54,
            0x78,
            0x9C,
            0x63,
            0x00,
            0x01,
            0x00,
            0x00,
            0x05,
            0x00,
            0x01,
            0x0D,
            0x0A,
            0x2D,
            0xB4,
            0x00,
            0x00,
            0x00,
            0x00,
            0x49,
            0x45,
            0x4E,
            0x44,
            0xAE,  # IEND chunk
            0x42,
            0x60,
            0x82,
        ]
    )
    test_image = tmp_path / "test.png"
    test_image.write_bytes(png_data)

    result = _image_to_base64(test_image)
    assert result.startswith("data:image/png;base64,")


def test_image_to_base64_gif_mime_type(tmp_path):
    """Should use correct MIME type for GIF."""
    # Create minimal GIF
    gif_data = b"GIF89a\x01\x00\x01\x00\x00\x00\x00;\x00\x00\x00\x00\x00\x00"
    test_image = tmp_path / "test.gif"
    test_image.write_bytes(gif_data)

    result = _image_to_base64(test_image)
    assert result.startswith("data:image/gif;base64,")


def test_create_slide_with_content():
    """Should create slide with title and bullets."""
    prs = Presentation()
    content = {
        "title": "Test Title",
        "subtitle": "Test Subtitle",
        "bullets": ["Point 1", "Point 2"],
        "notes": "Test speaker notes",
    }

    _create_slide(prs, content)

    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert slide.shapes.title.text == "Test Title"


def test_create_slide_empty_content():
    """Should handle empty content gracefully."""
    prs = Presentation()
    content = {"title": "", "subtitle": "", "bullets": [], "notes": ""}

    _create_slide(prs, content)

    assert len(prs.slides) == 1


# --- Integration tests with mocked OpenAI ---


def _create_test_png(path: Path) -> None:
    """Create a minimal valid PNG file for testing."""
    png_data = bytes(
        [
            0x89,
            0x50,
            0x4E,
            0x47,
            0x0D,
            0x0A,
            0x1A,
            0x0A,  # PNG signature
            0x00,
            0x00,
            0x00,
            0x0D,
            0x49,
            0x48,
            0x44,
            0x52,  # IHDR chunk
            0x00,
            0x00,
            0x00,
            0x01,
            0x00,
            0x00,
            0x00,
            0x01,
            0x08,
            0x06,
            0x00,
            0x00,
            0x00,
            0x1F,
            0x15,
            0xC4,
            0x89,
            0x00,
            0x00,
            0x00,
            0x0A,
            0x49,
            0x44,
            0x41,  # IDAT chunk
            0x54,
            0x78,
            0x9C,
            0x63,
            0x00,
            0x01,
            0x00,
            0x00,
            0x05,
            0x00,
            0x01,
            0x0D,
            0x0A,
            0x2D,
            0xB4,
            0x00,
            0x00,
            0x00,
            0x00,
            0x49,
            0x45,
            0x4E,
            0x44,
            0xAE,  # IEND chunk
            0x42,
            0x60,
            0x82,
        ]
    )
    path.write_bytes(png_data)


@pytest.fixture
def mock_openai_response():
    """Create a mock OpenAI API response."""
    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = json.dumps(
        {
            "title": "Test Slide Title",
            "subtitle": "Test Subtitle",
            "bullets": ["Point 1", "Point 2", "Point 3"],
            "notes": "These are speaker notes",
        }
    )
    return mock_response


def test_extract_slide_content_with_mock(tmp_path, mock_openai_response):
    """Should extract slide content using OpenAI API."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_client = MagicMock()
    mock_client.chat.completions.create.return_value = mock_openai_response

    result = _extract_slide_content(mock_client, test_image)

    assert result["title"] == "Test Slide Title"
    assert result["subtitle"] == "Test Subtitle"
    assert len(result["bullets"]) == 3
    assert result["notes"] == "These are speaker notes"
    mock_client.chat.completions.create.assert_called_once()


def test_extract_slide_content_api_error(tmp_path):
    """Should raise ValueError on generic OpenAI API error."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_client = MagicMock()
    mock_client.chat.completions.create.side_effect = OpenAIError("Connection failed")

    with pytest.raises(ValueError, match="OpenAI API error"):
        _extract_slide_content(mock_client, test_image)


def test_extract_slide_content_rate_limit_error(tmp_path):
    """Should raise ValueError with specific message on rate limit."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_client = MagicMock()
    mock_client.chat.completions.create.side_effect = RateLimitError(
        "Rate limit exceeded", response=MagicMock(), body=None
    )

    with pytest.raises(ValueError, match="rate limit exceeded"):
        _extract_slide_content(mock_client, test_image)


def test_extract_slide_content_auth_error(tmp_path):
    """Should raise ValueError with specific message on auth error."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_client = MagicMock()
    mock_client.chat.completions.create.side_effect = AuthenticationError(
        "Invalid API key", response=MagicMock(), body=None
    )

    with pytest.raises(ValueError, match="Invalid OpenAI API key"):
        _extract_slide_content(mock_client, test_image)


def test_extract_slide_content_malformed_json(tmp_path):
    """Should raise ValueError on malformed JSON response."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = "not valid json {"

    mock_client = MagicMock()
    mock_client.chat.completions.create.return_value = mock_response

    with pytest.raises(ValueError, match="Failed to parse GPT response as JSON"):
        _extract_slide_content(mock_client, test_image)


def test_images_to_pptx_impl_integration(tmp_path, mock_openai_response):
    """Integration test for full image-to-PPTX flow."""
    # Create test images
    image1 = tmp_path / "slide1.png"
    image2 = tmp_path / "slide2.png"
    _create_test_png(image1)
    _create_test_png(image2)

    output_path = tmp_path / "output.pptx"

    with patch("img2pptx.tools._get_client") as mock_get_client:
        mock_client = MagicMock()
        mock_client.chat.completions.create.return_value = mock_openai_response
        mock_get_client.return_value = mock_client

        result = _images_to_pptx_impl([str(image1), str(image2)], str(output_path))

    assert "Created PPTX with 2 slides" in result
    assert output_path.exists()

    # Verify the PPTX content
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 2


def test_images_to_pptx_validates_all_before_api_call(tmp_path, mock_openai_response):
    """Should validate all images before making any API calls."""
    valid_image = tmp_path / "valid.png"
    _create_test_png(valid_image)
    invalid_image = tmp_path / "missing.png"  # Does not exist

    output_path = tmp_path / "output.pptx"

    with patch("img2pptx.tools._get_client") as mock_get_client:
        mock_client = MagicMock()
        mock_client.chat.completions.create.return_value = mock_openai_response
        mock_get_client.return_value = mock_client

        result = _images_to_pptx_impl([str(valid_image), str(invalid_image)], str(output_path))

    assert "Error" in result
    # API should NOT have been called since validation failed
    mock_client.chat.completions.create.assert_not_called()


def test_extract_slide_content_empty_response(tmp_path):
    """Should raise ValueError on empty API response."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = None

    mock_client = MagicMock()
    mock_client.chat.completions.create.return_value = mock_response

    with pytest.raises(ValueError, match="API returned empty response"):
        _extract_slide_content(mock_client, test_image)


def test_extract_slide_content_missing_fields(tmp_path):
    """Should raise ValueError when response is missing required fields."""
    test_image = tmp_path / "slide.png"
    _create_test_png(test_image)

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = json.dumps(
        {
            "title": "Test Title",
            # Missing: subtitle, bullets, notes
        }
    )

    mock_client = MagicMock()
    mock_client.chat.completions.create.return_value = mock_response

    with pytest.raises(ValueError, match="missing required fields"):
        _extract_slide_content(mock_client, test_image)
