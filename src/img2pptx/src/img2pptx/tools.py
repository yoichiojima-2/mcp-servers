"""
Convert images to PPTX using OpenAI GPT-5.2.

Takes image slides (e.g., from nano-banana) and converts them to proper
editable PowerPoint presentations by using GPT-5.2's vision capabilities
to extract content and structure.
"""

import base64
import json
import os
from pathlib import Path

from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt

from core import SHARED_WORKSPACE, get_workspace

from . import mcp

# Constants
MAX_IMAGE_SIZE = 20_000_000  # 20MB
ALLOWED_EXTENSIONS = frozenset([".png", ".jpg", ".jpeg", ".webp", ".gif"])
FORBIDDEN_PATHS = frozenset(
    [
        "/bin",
        "/sbin",
        "/usr",
        "/etc",
        "/sys",
        "/proc",
        "/root",
        "/private/bin",
        "/private/etc",  # macOS symlinks
    ]
)


def _get_client() -> OpenAI:
    """Get configured OpenAI client."""
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OPENAI_API_KEY environment variable required")
    return OpenAI(api_key=api_key)


def _validate_image_path(path: Path) -> None:
    """Validate input image path."""
    if not path.exists():
        raise FileNotFoundError(f"Image not found: {path}")
    if path.stat().st_size > MAX_IMAGE_SIZE:
        raise ValueError(f"Image too large (max {MAX_IMAGE_SIZE // 1_000_000}MB)")
    if path.suffix.lower() not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported format: {path.suffix}")


def _validate_output_path(path: Path) -> None:
    """Validate output path is safe to write to."""
    resolved = str(path.resolve())
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Output must be .pptx, got: {path.suffix}")
    for forbidden in FORBIDDEN_PATHS:
        if resolved.startswith(forbidden + "/"):
            raise ValueError(f"Cannot write to system directory: {forbidden}")
    if not path.parent.exists():
        path.parent.mkdir(parents=True, exist_ok=True)


def _image_to_base64(image_path: Path) -> str:
    """Convert image to base64 data URL."""
    mime_types = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".gif": "image/gif",
    }
    mime = mime_types.get(image_path.suffix.lower(), "image/png")
    with open(image_path, "rb") as f:
        encoded = base64.b64encode(f.read()).decode("utf-8")
    return f"data:{mime};base64,{encoded}"


def _extract_slide_content(client: OpenAI, image_path: Path) -> dict:
    """Use GPT-5.2 to extract slide content from image."""
    base64_image = _image_to_base64(image_path)

    response = client.chat.completions.create(
        model="gpt-5.2",
        messages=[
            {
                "role": "system",
                "content": """You are a slide content extractor. Analyze the slide image and extract:
1. title: The main heading/title text
2. subtitle: Any subtitle or tagline (empty string if none)
3. bullets: List of bullet points (empty list if none)
4. notes: Any additional context or speaker notes you can infer

Respond in JSON format:
{"title": "...", "subtitle": "...", "bullets": ["...", "..."], "notes": "..."}""",
            },
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract the content from this slide image:"},
                    {"type": "image_url", "image_url": {"url": base64_image}},
                ],
            },
        ],
        response_format={"type": "json_object"},
    )

    return json.loads(response.choices[0].message.content)


def _create_slide(prs: Presentation, content: dict) -> None:
    """Create a slide from extracted content."""
    # Use title and content layout, with fallback
    try:
        layout = prs.slide_layouts[1]  # Title and Content
    except IndexError:
        layout = prs.slide_layouts[0]  # Fallback to first available
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = content.get("title", "")

    # Set subtitle/bullets in body
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Body placeholder
            tf = shape.text_frame
            tf.clear()

            subtitle = content.get("subtitle", "")
            if subtitle:
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(18)
                p.font.bold = True

            for bullet in content.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    # Add speaker notes
    notes = content.get("notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def _images_to_pptx_impl(image_paths: list[str], output_path: str) -> str:
    """Convert multiple images to a single PPTX."""
    if not image_paths:
        return "Error: No images provided"

    try:
        client = _get_client()
        output = Path(output_path).expanduser().resolve()
        _validate_output_path(output)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for img_path in image_paths:
            path = Path(img_path).expanduser().resolve()
            _validate_image_path(path)
            content = _extract_slide_content(client, path)
            _create_slide(prs, content)

        prs.save(output)
        return f"Created PPTX with {len(image_paths)} slides: {output}"

    except Exception as e:
        return f"Error: {e}"


def _image_to_pptx_impl(image_path: str, output_path: str) -> str:
    """Convert a single image to PPTX."""
    return _images_to_pptx_impl([image_path], output_path)


@mcp.tool()
def get_workspace_path() -> str:
    """Get the workspace directory path for saving PPTX files.

    Returns:
        Path to ~/.mcp-servers/workspace/
    """
    return str(get_workspace(SHARED_WORKSPACE))


@mcp.tool()
def image_to_pptx(image_path: str, output_path: str) -> str:
    """
    Convert a slide image to an editable PPTX presentation.

    Uses OpenAI GPT-5.2 vision to extract text and structure from the image,
    then creates a proper PowerPoint with editable text.

    Args:
        image_path: Path to the slide image (e.g., from nano-banana)
        output_path: Path for the output PPTX file

    Returns:
        Path to the created PPTX file, or error message

    Example:
        image_to_pptx(
            image_path="~/.mcp-servers/workspace/slide1.png",
            output_path="~/.mcp-servers/workspace/presentation.pptx"
        )
    """
    return _image_to_pptx_impl(image_path, output_path)


@mcp.tool()
def images_to_pptx(image_paths: list[str], output_path: str) -> str:
    """
    Convert multiple slide images to a single PPTX presentation.

    Uses OpenAI GPT-5.2 vision to extract text and structure from each image,
    creating a multi-slide PowerPoint with editable content.

    Args:
        image_paths: List of paths to slide images (in order)
        output_path: Path for the output PPTX file

    Returns:
        Path to the created PPTX file, or error message

    Example:
        images_to_pptx(
            image_paths=[
                "~/.mcp-servers/workspace/title.png",
                "~/.mcp-servers/workspace/content1.png",
                "~/.mcp-servers/workspace/content2.png"
            ],
            output_path="~/.mcp-servers/workspace/deck.pptx"
        )
    """
    return _images_to_pptx_impl(image_paths, output_path)
